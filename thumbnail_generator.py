#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPTXアイキャッチ画像生成クラス
"""

from pptx import Presentation
from pptx.util import Pt
import os
import subprocess


class ThumbnailGenerator:
    """PPTXテンプレートからアイキャッチ画像を生成するクラス"""
    
    def __init__(self, template_path: str):
        """
        初期化
        
        Args:
            template_path: PPTXテンプレートファイルのパス
        """
        self.template_path = template_path
    
    def create_thumbnail(self, title: str, output_path: str) -> bool:
        """
        タイトルを挿入してアイキャッチ画像を生成
        
        Args:
            title: 記事タイトル
            output_path: 出力画像のパス（.png）
            
        Returns:
            成功時True、失敗時False
        """
        try:
            # PPTXを読み込み
            prs = Presentation(self.template_path)
            
            # 最初のスライドを取得
            slide = prs.slides[0]
            
            # タイトルテキストを探して置き換え
            # 通常、2番目のテキストボックスがタイトル部分
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # "私たちの理念と風土" のようなテキストを探す
                    if "私たちの理念" in shape.text or len(shape.text) > 5:
                        # テキストを置き換え
                        text_frame = shape.text_frame
                        text_frame.clear()
                        p = text_frame.paragraphs[0]
                        run = p.add_run()
                        run.text = title
                        
                        # フォントサイズを調整（タイトルの長さに応じて）
                        if len(title) > 30:
                            run.font.size = Pt(32)
                        elif len(title) > 20:
                            run.font.size = Pt(40)
                        else:
                            run.font.size = Pt(48)
                        
                        break
            
            # 一時的なPPTXファイルとして保存
            temp_pptx = output_path.replace('.png', '_temp.pptx')
            prs.save(temp_pptx)
            
            # LibreOfficeを使ってPNGに変換
            output_dir = os.path.dirname(output_path)
            result = subprocess.run(
                ['libreoffice', '--headless', '--convert-to', 'png', 
                 '--outdir', output_dir, temp_pptx],
                capture_output=True,
                text=True,
                timeout=30
            )
            
            # 生成されたPNGファイル名を変更
            temp_png = temp_pptx.replace('.pptx', '.png')
            if os.path.exists(temp_png):
                os.rename(temp_png, output_path)
            
            # 一時ファイルを削除
            if os.path.exists(temp_pptx):
                os.remove(temp_pptx)
            
            return os.path.exists(output_path)
        
        except Exception as e:
            print(f"サムネイル生成エラー: {e}")
            return False


if __name__ == '__main__':
    # テスト用コード
    generator = ThumbnailGenerator('/home/ubuntu/upload/ブルーシンプル資料noteアイキャッチ.pptx')
    success = generator.create_thumbnail(
        'ショーペンハウアー｜「人生は振り子」から考える、満たされない心の正体',
        '/home/ubuntu/test_thumbnail.png'
    )
    print(f"サムネイル生成: {'成功' if success else '失敗'}")
